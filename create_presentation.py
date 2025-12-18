"""
Generate PowerPoint presentation for Home Price Prediction project.
Uses Georgia font, 30pt headlines, 16pt body text.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import json

# Constants
GEORGIA = "Georgia"
HEADLINE_SIZE = Pt(30)
BODY_SIZE = Pt(16)
SUBHEAD_SIZE = Pt(20)

# Colors
PRIMARY_COLOR = RGBColor(102, 126, 234)  # #667eea
SECONDARY_COLOR = RGBColor(118, 75, 162)  # #764ba2
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(45, 55, 72)  # #2d3748
SUCCESS_COLOR = RGBColor(72, 187, 120)  # Green

def add_title_slide(prs, title, subtitle):
    """Add a title slide with gradient-style header"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = GEORGIA
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.name = GEORGIA
    p.font.size = SUBHEAD_SIZE
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets):
    """Add a content slide with headline and bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = GEORGIA
    p.font.size = HEADLINE_SIZE
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = GEORGIA
        p.font.size = BODY_SIZE
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(12)
    
    return slide

def add_results_slide(prs, title, metrics):
    """Add a results slide with key metrics"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = GEORGIA
    p.font.size = HEADLINE_SIZE
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Metrics boxes
    y_start = 1.8
    x_positions = [0.5, 3.5, 6.5]
    
    for i, (label, value) in enumerate(metrics[:3]):
        x = x_positions[i]
        
        # Metric box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y_start), Inches(2.8), Inches(1.8)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(248, 250, 252)
        box.line.color.rgb = PRIMARY_COLOR
        
        # Value
        val_box = slide.shapes.add_textbox(Inches(x), Inches(y_start + 0.3), Inches(2.8), Inches(0.8))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(value)
        p.font.name = GEORGIA
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        lbl_box = slide.shapes.add_textbox(Inches(x), Inches(y_start + 1.1), Inches(2.8), Inches(0.5))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.name = GEORGIA
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER
    
    # Additional metrics below
    if len(metrics) > 3:
        y_start = 4.0
        for i, (label, value) in enumerate(metrics[3:6]):
            if i >= 3:
                break
            x = x_positions[i]
            
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y_start), Inches(2.8), Inches(1.5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(248, 250, 252)
            box.line.color.rgb = SECONDARY_COLOR
            
            val_box = slide.shapes.add_textbox(Inches(x), Inches(y_start + 0.2), Inches(2.8), Inches(0.6))
            tf = val_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(value)
            p.font.name = GEORGIA
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = SECONDARY_COLOR
            p.alignment = PP_ALIGN.CENTER
            
            lbl_box = slide.shapes.add_textbox(Inches(x), Inches(y_start + 0.85), Inches(2.8), Inches(0.5))
            tf = lbl_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.name = GEORGIA
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_TEXT
            p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Load metrics from saved results
    ROOT = Path(__file__).parent
    MODELS_DIR = ROOT / 'models'
    
    # Try to load ensemble summary (new format)
    ensemble_summary = {}
    ensemble_path = MODELS_DIR / 'ensemble_summary.json'
    if ensemble_path.exists():
        with open(ensemble_path) as f:
            ensemble_summary = json.load(f)
    
    # Try to load XGBoost metrics (fallback)
    xgb_metrics = {}
    xgb_path = MODELS_DIR / 'xgboost_enhanced_metrics.json'
    if xgb_path.exists():
        with open(xgb_path) as f:
            xgb_metrics = json.load(f)
    
    baseline_summary = {}
    baseline_path = MODELS_DIR / 'baseline_models_summary.json'
    if baseline_path.exists():
        with open(baseline_path) as f:
            baseline_summary = json.load(f)
    
    # Use ensemble results if available, fallback to XGBoost
    if ensemble_summary:
        best_r2 = ensemble_summary.get('test_r2', 0.8561)
        best_rmse = ensemble_summary.get('rmse', 321335)
        best_model = ensemble_summary.get('best_model', 'XGB+LGB Blend')
    else:
        best_r2 = xgb_metrics.get('test_metrics', {}).get('r2', 0.8506)
        best_rmse = xgb_metrics.get('test_metrics', {}).get('rmse', 327396)
        best_model = 'XGBoost GPU-Tuned'
    
    # === SLIDE 1: Title ===
    add_title_slide(prs, 
        "Home Price Prediction",
        "Machine Learning for Real Estate Valuation\nRutgers University • December 2025"
    )
    
    # === SLIDE 2: Project Overview ===
    add_content_slide(prs, "Project Overview", [
        "Objective: Predict home prices with high accuracy (targeting 88.4% R²)",
        "Data: Louisiana MLS real estate transactions (January - August 2025)",
        "173,070 total properties • 150,311 training • 22,759 testing",
        "1,022 engineered features including location, property details, and amenities",
        "Chronological train/test split (months 1-7 / month 8)",
        "Benchmark to beat: Steph's Random Forest at 88.4% R²"
    ])
    
    # === SLIDE 3: Methodology ===
    add_content_slide(prs, "Methodology", [
        "Data Pipeline: Load → Clean → Feature Engineer → Train → Evaluate",
        "Feature Engineering: Sanitized names, building age, geographic encoding",
        "Leakage Prevention: Removed ListPrice, CloseDate, agent/office info",
        "Baseline Models: Linear Regression, Ridge, Lasso, Random Forest",
        "Advanced Models: XGBoost (GPU), LightGBM, CatBoost, Neural Networks",
        "Ensemble Methods: Voting, Stacking, Weighted Blending",
        "Hyperparameter Tuning: 3-stage RandomizedSearchCV with GPU acceleration"
    ])
    
    # === SLIDE 4: XGBoost Architecture ===
    add_content_slide(prs, "XGBoost & LightGBM GPU-Accelerated Tuning", [
        "XGBoost: 3-stage tuning (broad → refined → final training)",
        "LightGBM: RandomizedSearchCV with 20 iterations, 3-fold CV",
        "GPU Acceleration: device='cuda' for XGBoost on Amarel cluster",
        "LightGBM: CPU mode with parallel processing (n_jobs=-1)",
        "Key Parameters: learning_rate, num_leaves, max_depth, colsample",
        "Regularization: reg_alpha (L1), reg_lambda (L2) for both models",
        "Early Stopping: 50 rounds patience to prevent overfitting"
    ])
    
    # === SLIDE 5: Results ===
    add_results_slide(prs, "Model Performance Results", [
        ("Best R² Score", f"{best_r2*100:.2f}%"),
        ("RMSE", f"${best_rmse:,.0f}"),
        ("Best Model", best_model.replace('_', ' ')),
        ("Training Samples", "150,311"),
        ("Test Samples", "22,759"),
        ("Features Used", "1,022"),
    ])
    
    # === SLIDE 6: Model Comparison ===
    add_content_slide(prs, "Model Comparison", [
        f"Weighted Blend (10% XGB + 90% LGB): 85.61% R² ← BEST",
        "LightGBM (Tuned): 85.60% R²",
        "XGBoost (GPU-Tuned): 85.06% R²",
        "Stacking (XGB + LGB + Ridge meta): 83.95% R²",
        "Voting (XGB + LGB average): 83.25% R²",
        "Random Forest Baseline: 88.4% R² (Steph's benchmark)",
        f"Gap to Steph: -2.79 percentage points (room for improvement)"
    ])
    
    # === SLIDE 7: Key Features ===
    add_content_slide(prs, "Top Predictive Features", [
        "Living Area (sq ft) - Most important predictor",
        "Building Age - Derived from YearBuilt",
        "Geographic Location - City, MLSAreaMajor, Postal Code (target encoded)",
        "Bedrooms & Bathrooms - Core property attributes",
        "Garage Spaces - Storage and parking value",
        "Property Type - Single Family, Condo, Multi-Family",
        "Property Condition & Quality - Assessor ratings"
    ])
    
    # === SLIDE 8: Technical Stack ===
    add_content_slide(prs, "Technical Implementation", [
        "Python 3.9+ with scikit-learn, XGBoost 2.1, LightGBM 3.3",
        "GPU: NVIDIA CUDA via device='cuda' on Amarel HPC cluster",
        "Data Processing: pandas, numpy for 173K+ records × 1K features",
        "Feature Names: Sanitized for LightGBM JSON compatibility",
        "Web App: Streamlit with plotly for interactive predictions",
        "Notebooks: Jupyter for reproducible analysis pipeline",
        "SLURM: Job submission for HPC GPU computing"
    ])
    
    # === SLIDE 9: Challenges & Solutions ===
    add_content_slide(prs, "Challenges & Solutions", [
        "Challenge: Large dataset (150K records × 1,022 features)",
        "Solution: GPU acceleration + efficient memory management",
        "Challenge: Feature name compatibility (LightGBM JSON limits)",
        "Solution: Sanitize names (replace commas, spaces with underscores)",
        "Challenge: LightGBM GPU build not available on cluster",
        "Solution: Use CPU mode with n_jobs=-1 (still fast: 21s training)",
        "Challenge: Gap to Steph's 88.4% baseline",
        "Solution: Ensemble blending brought us from 84.68% to 85.61%"
    ])
    
    # === SLIDE 10: Future Work ===
    add_content_slide(prs, "Future Improvements", [
        "Feature Engineering: More geographic/neighborhood features",
        "Ensemble Optimization: Grid search over blend weights",
        "Deep Learning: Neural networks for complex patterns",
        "CatBoost: Not tested yet - native categorical handling",
        "Time Series: Add temporal market trend features",
        "External Data: Economic indicators, interest rates, school ratings",
        "Target: Close the 2.79% gap to Steph's 88.4% baseline"
    ])
    
    # === SLIDE 11: Conclusion ===
    add_content_slide(prs, "Conclusion", [
        f"Achieved {best_r2*100:.2f}% R² on test set with {best_model}",
        "Best result: Weighted Blend (10% XGBoost + 90% LightGBM)",
        f"RMSE of ${best_rmse:,.0f} - average prediction error",
        "Gap to Steph's 88.4%: 2.79 percentage points remaining",
        "Clean, reproducible pipeline validated on Amarel HPC",
        "Web application ready for interactive predictions",
        "Strong foundation for continued improvement"
    ])
    
    # === SLIDE 12: Thank You ===
    add_title_slide(prs,
        "Thank You!",
        "Questions?\n\nHome Price Prediction Project • Rutgers 2025"
    )
    
    # Save presentation
    output_path = ROOT / 'Home_Price_Prediction_Presentation.pptx'
    prs.save(output_path)
    print(f"✓ Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
